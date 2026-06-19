"""
Module 0 — Structural Baseline Establishment.

Analyzes document structure to classify text runs into roles (header, table, body, etc.)
and establishes expected font/size patterns PER ROLE before any anomaly detection.

This baseline is consumed by downstream modules (fingerprint, OCR, heatmap) to avoid
false positives on legitimate structural variations like bold headers or table labels.
"""

import statistics
from collections import defaultdict, Counter
from typing import Dict, List, Tuple, Any

from pptx import Presentation
import fitz  # PyMuPDF
from docx import Document


def establish_baseline(file_path: str, file_type: str) -> dict:
    """
    Establish structural baseline for document analysis.
    
    Args:
        file_path: Path to the file
        file_type: 'pptx', 'pdf', or 'docx'
    
    Returns:
        Dict with role classifications and per-role font/size expectations
    """
    if file_type == "pptx":
        return _establish_pptx_baseline(file_path)
    elif file_type == "pdf":
        return _establish_pdf_baseline(file_path)
    elif file_type == "docx":
        return _establish_docx_baseline(file_path)
    else:
        return {
            "roles": {},
            "role_expectations": {},
            "total_runs": 0
        }


def _establish_pdf_baseline(file_path: str) -> dict:
    """Establish structural baseline for PDF files."""
    doc = fitz.open(file_path)
    
    all_runs = []  # All text spans with metadata
    page_dimensions = []
    
    # First pass: collect all runs with position and styling
    for page_idx in range(len(doc)):
        page = doc[page_idx]
        page_rect = page.rect
        page_dimensions.append({
            "width": page_rect.width,
            "height": page_rect.height
        })
        
        blocks = page.get_text("dict")["blocks"]
        
        for block in blocks:
            if block.get("type") != 0:  # Skip non-text blocks
                continue
            
            bbox = block.get("bbox", [0, 0, 0, 0])
            
            for line in block.get("lines", []):
                for span in line.get("spans", []):
                    text = span.get("text", "").strip()
                    if not text:
                        continue
                    
                    font_name = span.get("font", "Unknown")
                    if "+" in font_name:
                        font_name = font_name.split("+", 1)[1]
                    
                    size = span.get("size", 12)
                    flags = span.get("flags", 0)
                    bold = bool(flags & 16)
                    italic = bool(flags & 2)
                    
                    all_runs.append({
                        "page": page_idx + 1,
                        "text": text,
                        "font": font_name,
                        "size": size,
                        "bold": bold,
                        "italic": italic,
                        "x": bbox[0],
                        "y": bbox[1],
                        "width": bbox[2] - bbox[0],
                        "height": bbox[3] - bbox[1],
                        "page_height": page_rect.height,
                        "page_width": page_rect.width
                    })
    
    doc.close()
    
    if not all_runs:
        return {
            "roles": {},
            "role_expectations": {},
            "total_runs": 0
        }
    
    # Calculate global statistics
    all_sizes = [r["size"] for r in all_runs]
    median_size = statistics.median(all_sizes)
    
    # Detect tables (simplified: aligned vertical/horizontal text clusters)
    table_regions = _detect_table_regions(all_runs)
    
    # Classify each run into a structural role
    role_assignments = {}
    for idx, run in enumerate(all_runs):
        role = _classify_run_role(run, median_size, table_regions)
        role_assignments[idx] = role
        run["role"] = role
    
    # Compute per-role expectations
    role_expectations = _compute_role_expectations(all_runs)
    
    # Group runs by role for reporting
    roles_summary = defaultdict(int)
    for role in role_assignments.values():
        roles_summary[role] += 1
    
    return {
        "roles": dict(roles_summary),
        "role_expectations": role_expectations,
        "total_runs": len(all_runs),
        "median_size": median_size,
        "run_details": all_runs  # Full data for downstream modules
    }


def _establish_pptx_baseline(file_path: str) -> dict:
    """Establish structural baseline for PPTX files."""
    prs = Presentation(file_path)
    
    all_runs = []
    
    for slide_idx, slide in enumerate(prs.slides):
        slide_width = int(prs.slide_width / 914400 * 96)
        slide_height = int(prs.slide_height / 914400 * 96)
        
        for shape in slide.shapes:
            if not hasattr(shape, "text_frame"):
                continue
            
            shape_left = int(shape.left / 914400 * 96) if hasattr(shape, "left") else 0
            shape_top = int(shape.top / 914400 * 96) if hasattr(shape, "top") else 0
            
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text = run.text.strip()
                    if not text:
                        continue
                    
                    font = run.font
                    font_name = font.name if font.name else "Default"
                    size = int(font.size.pt) if font.size else 12
                    bold = font.bold if font.bold is not None else False
                    italic = font.italic if font.italic is not None else False
                    
                    all_runs.append({
                        "page": slide_idx + 1,
                        "text": text,
                        "font": font_name,
                        "size": size,
                        "bold": bold,
                        "italic": italic,
                        "x": shape_left,
                        "y": shape_top,
                        "page_height": slide_height,
                        "page_width": slide_width
                    })
    
    if not all_runs:
        return {
            "roles": {},
            "role_expectations": {},
            "total_runs": 0
        }
    
    all_sizes = [r["size"] for r in all_runs]
    median_size = statistics.median(all_sizes)
    
    table_regions = []  # PPTX table detection would need more complex logic
    
    role_assignments = {}
    for idx, run in enumerate(all_runs):
        role = _classify_run_role(run, median_size, table_regions)
        role_assignments[idx] = role
        run["role"] = role
    
    role_expectations = _compute_role_expectations(all_runs)
    
    roles_summary = defaultdict(int)
    for role in role_assignments.values():
        roles_summary[role] += 1
    
    return {
        "roles": dict(roles_summary),
        "role_expectations": role_expectations,
        "total_runs": len(all_runs),
        "median_size": median_size,
        "run_details": all_runs
    }


def _establish_docx_baseline(file_path: str) -> dict:
    """Establish structural baseline for DOCX files."""
    doc = Document(file_path)
    
    all_runs = []
    
    for para_idx, paragraph in enumerate(doc.paragraphs):
        for run in paragraph.runs:
            text = run.text.strip()
            if not text:
                continue
            
            font = run.font
            font_name = font.name if font.name else "Default"
            size = int(font.size.pt) if font.size else 12
            bold = font.bold if font.bold is not None else False
            italic = font.italic if font.italic is not None else False
            
            all_runs.append({
                "page": para_idx // 20 + 1,  # Approximate page
                "text": text,
                "font": font_name,
                "size": size,
                "bold": bold,
                "italic": italic,
                "para_idx": para_idx
            })
    
    if not all_runs:
        return {
            "roles": {},
            "role_expectations": {},
            "total_runs": 0
        }
    
    all_sizes = [r["size"] for r in all_runs]
    median_size = statistics.median(all_sizes)
    
    table_regions = []
    
    role_assignments = {}
    for idx, run in enumerate(all_runs):
        role = _classify_run_role(run, median_size, table_regions)
        role_assignments[idx] = role
        run["role"] = role
    
    role_expectations = _compute_role_expectations(all_runs)
    
    roles_summary = defaultdict(int)
    for role in role_assignments.values():
        roles_summary[role] += 1
    
    return {
        "roles": dict(roles_summary),
        "role_expectations": role_expectations,
        "total_runs": len(all_runs),
        "median_size": median_size,
        "run_details": all_runs
    }


def _classify_run_role(run: dict, median_size: float, table_regions: List[Tuple]) -> str:
    """
    Classify a text run into a structural role based on position and styling.
    
    Args:
        run: Dict with text, font, size, bold, position info
        median_size: Median font size across document
        table_regions: List of (x1, y1, x2, y2) bounding boxes for detected tables
    
    Returns:
        Role string: "header", "table_cell", "table_header", "body_text", "label", "footer"
    """
    text = run.get("text", "")
    size = run.get("size", median_size)
    bold = run.get("bold", False)
    y_pos = run.get("y", 0)
    page_height = run.get("page_height", 1000)
    
    # Check if inside a table
    x = run.get("x", 0)
    for x1, y1, x2, y2 in table_regions:
        if x1 <= x <= x2 and y1 <= y_pos <= y2:
            # Inside table - check if header row (bold + top of table)
            if bold and abs(y_pos - y1) < 30:
                return "table_header"
            return "table_cell"
    
    # Header/Title: Large size OR (bold + top 15% of page)
    if size > median_size * 1.4:
        return "header"
    
    if bold and size >= median_size * 1.1 and y_pos < page_height * 0.15:
        return "header"
    
    # Footer: Bottom 10% of page OR very small text
    if y_pos > page_height * 0.9 or size < median_size * 0.75:
        return "footer"
    
    # Label: Short text (< 20 chars) + ends with colon OR bold + short
    if len(text) < 20 and (text.endswith(":") or (bold and len(text) < 15)):
        return "label"
    
    # Default: body text
    return "body_text"


def _detect_table_regions(runs: List[dict]) -> List[Tuple]:
    """
    Detect table regions based on aligned text clusters.
    Simplified heuristic: groups of runs with similar x-coordinates (columns).
    
    Returns:
        List of (x1, y1, x2, y2) bounding boxes
    """
    if not runs:
        return []
    
    # Group runs by approximate x-coordinate (quantized to 20px buckets)
    x_buckets = defaultdict(list)
    for run in runs:
        x = run.get("x", 0)
        bucket = int(x / 20) * 20
        x_buckets[bucket].append(run)
    
    # If we have 3+ vertical columns with 5+ items each, likely a table
    dense_columns = [runs for runs in x_buckets.values() if len(runs) >= 5]
    
    if len(dense_columns) >= 3:
        # Find bounding box of these columns
        all_x = [r["x"] for col in dense_columns for r in col]
        all_y = [r["y"] for col in dense_columns for r in col]
        
        if all_x and all_y:
            return [(min(all_x), min(all_y), max(all_x) + 200, max(all_y) + 50)]
    
    return []


def _compute_role_expectations(runs: List[dict]) -> dict:
    """
    Compute expected font/size patterns per structural role.
    
    Args:
        runs: List of run dicts with 'role', 'font', 'size', 'bold', etc.
    
    Returns:
        Dict mapping role -> expectations dict with font_range, size_range, common_styles
    """
    role_data = defaultdict(lambda: {
        "fonts": [],
        "sizes": [],
        "bold_count": 0,
        "italic_count": 0,
        "total": 0
    })
    
    for run in runs:
        role = run.get("role", "body_text")
        role_data[role]["fonts"].append(run.get("font", "Unknown"))
        role_data[role]["sizes"].append(run.get("size", 12))
        if run.get("bold"):
            role_data[role]["bold_count"] += 1
        if run.get("italic"):
            role_data[role]["italic_count"] += 1
        role_data[role]["total"] += 1
    
    expectations = {}
    
    for role, data in role_data.items():
        if data["total"] == 0:
            continue
        
        sizes = data["sizes"]
        fonts = data["fonts"]
        
        # Calculate size range (allow ±20% variation from median)
        median_size = statistics.median(sizes)
        size_min = median_size * 0.8
        size_max = median_size * 1.2
        
        # Most common fonts for this role
        font_counter = Counter(fonts)
        common_fonts = [f for f, _ in font_counter.most_common(3)]
        
        # Bold/italic prevalence
        bold_prevalence = data["bold_count"] / data["total"]
        italic_prevalence = data["italic_count"] / data["total"]
        
        expectations[role] = {
            "median_size": round(median_size, 1),
            "size_range": (round(size_min, 1), round(size_max, 1)),
            "common_fonts": common_fonts,
            "bold_expected": bold_prevalence > 0.5,  # >50% are bold
            "italic_expected": italic_prevalence > 0.3,
            "sample_count": data["total"],
            "tolerance": "high" if role in ["header", "label", "table_header"] else "medium"
        }
    
    return expectations

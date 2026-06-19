import io
import os
import subprocess
import tempfile
from typing import Dict, List, Any
from pptx import Presentation
import fitz  # PyMuPDF
from PIL import Image, ImageDraw
from docx import Document
from app.db.supabase_client import get_supabase_client


async def generate_tamper_heatmap(
    file_path: str, file_type: str, flagged_findings: dict
) -> dict:
    """
    Generate tamper detection heatmap for PPTX/PDF/DOCX files.
    
    Args:
        file_path: Path to the file
        file_type: 'pptx', 'pdf', or 'docx'
        flagged_findings: Dict containing findings from modules 1,2,3,5,6,7
    
    Returns:
        Dict with pages, image URLs, and bounding boxes with intensity levels
    """
    if file_type == "pptx":
        return await _process_pptx(file_path, flagged_findings)
    elif file_type == "pdf":
        return await _process_pdf(file_path, flagged_findings)
    elif file_type == "docx":
        return await _process_docx(file_path, flagged_findings)
    else:
        return {"pages": []}


async def _process_pptx(file_path: str, flagged_findings: dict) -> dict:
    """Process PPTX file and generate heatmap data."""
    prs = Presentation(file_path)
    pages = []
    
    # Extract document_id from findings if available
    doc_id = flagged_findings.get("document_id", "temp")
    
    for slide_idx, slide in enumerate(prs.slides):
        page_num = slide_idx + 1
        boxes = []
        
        # Extract all shapes and their bounding boxes
        for shape in slide.shapes:
            if hasattr(shape, "left") and hasattr(shape, "top"):
                # Convert EMU to pixels (914400 EMU = 1 inch, 96 DPI)
                x = int(shape.left / 914400 * 96)
                y = int(shape.top / 914400 * 96)
                width = int(shape.width / 914400 * 96)
                height = int(shape.height / 914400 * 96)
                
                # Extract text if available
                text_content = ""
                if hasattr(shape, "text"):
                    text_content = shape.text.strip()
                
                # Determine intensity based on flagged_findings
                intensity, reason, source = _match_finding(
                    page_num, text_content, flagged_findings
                )
                
                if intensity:
                    boxes.append({
                        "x": x,
                        "y": y,
                        "width": width,
                        "height": height,
                        "intensity": intensity,
                        "reason": reason,
                        "source_module": source
                    })
        
        # Render slide as image
        image_url = await _render_pptx_slide(file_path, slide_idx, doc_id)
        
        pages.append({
            "page_number": page_num,
            "image_url": image_url,
            "boxes": boxes
        })
    
    return {"pages": pages}


async def _process_pdf(file_path: str, flagged_findings: dict) -> dict:
    """Process PDF file and generate heatmap data."""
    doc = fitz.open(file_path)
    pages = []
    
    # Extract document_id from findings if available
    doc_id = flagged_findings.get("document_id", "temp")
    
    for page_idx in range(len(doc)):
        page = doc[page_idx]
        page_num = page_idx + 1
        boxes = []
        
        # Extract text and image blocks
        blocks = page.get_text("dict")["blocks"]
        
        for block in blocks:
            bbox = block.get("bbox")
            if not bbox:
                continue
            
            x, y, x1, y1 = bbox
            width = int(x1 - x)
            height = int(y1 - y)
            
            # Extract text content
            text_content = ""
            if block.get("type") == 0:  # Text block
                for line in block.get("lines", []):
                    for span in line.get("spans", []):
                        text_content += span.get("text", "") + " "
            
            text_content = text_content.strip()
            
            # Determine intensity based on flagged_findings
            intensity, reason, source = _match_finding(
                page_num, text_content, flagged_findings
            )
            
            if intensity:
                boxes.append({
                    "x": int(x),
                    "y": int(y),
                    "width": width,
                    "height": height,
                    "intensity": intensity,
                    "reason": reason,
                    "source_module": source
                })
        
        # Render page as image
        image_url = await _render_pdf_page(page, page_idx, doc_id)
        
        pages.append({
            "page_number": page_num,
            "image_url": image_url,
            "boxes": boxes
        })
    
    doc.close()
    return {"pages": pages}


def _match_finding(page_num: int, text_content: str, flagged_findings: dict) -> tuple:
    """
    Match a box against flagged findings to determine intensity.
    
    Returns:
        tuple: (intensity, reason, source_module) or (None, None, None)
    """
    # Check findings from all modules
    for module, findings in flagged_findings.items():
        if not isinstance(findings, dict):
            continue
        
        # Skip non-module keys
        if module in ["document_id"]:
            continue
        
        # Check if findings contain page/slide references
        if isinstance(findings, dict):
            # Page-level match
            if findings.get("page") == page_num or findings.get("slide") == page_num:
                return ("low", findings.get("message", "Flagged page"), module)
            
            # Text content match (high intensity)
            flagged_text = findings.get("flagged_text", "")
            if flagged_text and text_content and flagged_text.lower() in text_content.lower():
                return ("high", findings.get("message", "Flagged content"), module)
            
            # Check for anomalies lists (OCR, metadata, etc.)
            anomalies = findings.get("anomalies_found", []) or findings.get("font_anomalies", [])
            for anomaly in anomalies:
                if isinstance(anomaly, dict):
                    # Check location match
                    location = anomaly.get("location", "")
                    if f"Slide {page_num}" in location or f"Page {page_num}" in location:
                        return ("medium", anomaly.get("issue", "Anomaly detected"), module)
                    
                    # Check text match
                    anomaly_text = anomaly.get("text", "") or anomaly.get("expected", "")
                    if anomaly_text and text_content and anomaly_text.lower() in text_content.lower():
                        return ("high", anomaly.get("issue", "Flagged content"), module)
            
            # Check nested findings list
            if "findings" in findings and isinstance(findings["findings"], list):
                for finding in findings["findings"]:
                    if isinstance(finding, dict):
                        if finding.get("page") == page_num or finding.get("slide") == page_num:
                            return ("low", finding.get("message", "Flagged page"), module)
                        
                        finding_text = finding.get("text", "") or finding.get("content", "")
                        if finding_text and text_content and finding_text.lower() in text_content.lower():
                            return ("high", finding.get("message", "Flagged content"), module)
            
            # Check contradictions (module6)
            contradictions = findings.get("contradictions", [])
            for contradiction in contradictions:
                if isinstance(contradiction, dict):
                    claim = contradiction.get("claim", "")
                    if claim and text_content and claim.lower() in text_content.lower():
                        return ("high", "Contradictory statement", module)
            
            # Check compliance issues (module7)
            issues = findings.get("issues", [])
            for issue in issues:
                if isinstance(issue, dict):
                    issue_text = issue.get("text", "") or issue.get("description", "")
                    if issue_text and text_content and issue_text.lower() in text_content.lower():
                        return ("medium", issue.get("type", "Compliance issue"), module)
    
    return (None, None, None)


def _libreoffice_render(file_path: str, slide_idx: int, temp_dir: str) -> str | None:
    """
    Attempt a LibreOffice headless PNG conversion of *file_path*.
    Returns the path of the matching slide PNG, or None on any failure.
    Does NOT raise — callers treat None as "fall back to Pillow".
    """
    try:
        subprocess.run(
            ["soffice", "--headless", "--convert-to", "png", "--outdir", temp_dir, file_path],
            capture_output=True,
            timeout=30,
            check=False,
        )
        png_files = sorted(f for f in os.listdir(temp_dir) if f.endswith(".png"))
        if png_files and slide_idx < len(png_files):
            return os.path.join(temp_dir, png_files[slide_idx])
    except (FileNotFoundError, subprocess.TimeoutExpired, Exception):
        pass
    return None


def _pillow_render(file_path: str, slide_idx: int) -> Image.Image:
    """
    Render a PPTX slide as a Pillow Image using shape geometry and fill colours.
    This is a faithful placeholder: background colour is extracted from the slide
    theme when available; each shape is drawn with its fill colour (or a neutral
    grey) and a darker outline; text runs are rendered as single-line truncated
    labels so the spatial layout is recognisable.
    """
    from pptx.util import Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    EMU = 914400  # EMU per inch
    DPI = 96

    prs = Presentation(file_path)
    slide = prs.slides[slide_idx]

    W = max(1, int(prs.slide_width / EMU * DPI))
    H = max(1, int(prs.slide_height / EMU * DPI))

    # ── background colour ────────────────────────────────────────────────
    bg_rgb = (30, 32, 40)  # dark default matching app theme
    try:
        bg = slide.background
        fill = bg.fill
        if fill.type is not None and fill.fore_color.type is not None:
            c = fill.fore_color.rgb
            bg_rgb = (c.red, c.green, c.blue)
    except Exception:
        pass

    img = Image.new("RGB", (W, H), bg_rgb)
    draw = ImageDraw.Draw(img, "RGBA")

    for shape in slide.shapes:
        if not (hasattr(shape, "left") and shape.left is not None):
            continue

        x  = int(shape.left  / EMU * DPI)
        y  = int(shape.top   / EMU * DPI)
        bw = int(shape.width / EMU * DPI)
        bh = int(shape.height/ EMU * DPI)

        # ── shape fill ───────────────────────────────────────────────────
        fill_rgba = (255, 255, 255, 40)  # translucent white fallback
        outline_rgb = (120, 130, 150)
        try:
            sf = shape.fill
            if sf.type is not None:
                c = sf.fore_color.rgb
                fill_rgba = (c.red, c.green, c.blue, 80)
                outline_rgb = (max(0, c.red - 40), max(0, c.green - 40), max(0, c.blue - 40))
        except Exception:
            pass

        draw.rectangle([x, y, x + bw, y + bh], fill=fill_rgba, outline=outline_rgb, width=1)

        # ── text label ───────────────────────────────────────────────────
        try:
            raw = shape.text.strip()
            if raw:
                label = raw[:48] + ("…" if len(raw) > 48 else "")
                # centre the label inside the shape
                tx = x + 4
                ty = y + max(0, (bh - 10) // 2)
                draw.text((tx + 1, ty + 1), label, fill=(0, 0, 0, 120))   # shadow
                draw.text((tx,     ty    ), label, fill=(220, 220, 230, 220))
        except Exception:
            pass

    return img


async def _render_pptx_slide(file_path: str, slide_idx: int, doc_id: str) -> str:
    """Render PPTX slide as PNG and upload to Supabase.
    
    Tries LibreOffice headless first; falls back to the Pillow placeholder
    renderer on any failure.  Never raises — returns "" if upload also fails.
    """
    temp_dir = tempfile.mkdtemp()
    output_file = None
    used_libreoffice = False

    try:
        # ── 1. Try LibreOffice ────────────────────────────────────────────
        lo_file = _libreoffice_render(file_path, slide_idx, temp_dir)
        if lo_file:
            output_file = lo_file
            used_libreoffice = True

        # ── 2. Pillow fallback ────────────────────────────────────────────
        if not output_file:
            img = _pillow_render(file_path, slide_idx)
            output_file = os.path.join(temp_dir, f"slide_{slide_idx}.png")
            img.save(output_file, "PNG")

        # ── 3. Upload ─────────────────────────────────────────────────────
        return await _upload_to_supabase(output_file, doc_id, slide_idx)

    except Exception:
        return ""
    finally:
        # Always clean up temp dir regardless of success or failure
        try:
            for f in os.listdir(temp_dir):
                os.unlink(os.path.join(temp_dir, f))
            os.rmdir(temp_dir)
        except Exception:
            pass


async def _render_pdf_page(page, page_idx: int, doc_id: str) -> str:
    """Render PDF page as PNG and upload to Supabase."""
    try:
        # Render page to image using PyMuPDF
        pix = page.get_pixmap(dpi=150)
        img_data = pix.tobytes("png")
        
        # Save to temp file
        with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as tmp:
            tmp.write(img_data)
            tmp_path = tmp.name
        
        # Upload to Supabase
        url = await _upload_to_supabase(tmp_path, doc_id, page_idx)
        os.unlink(tmp_path)
        
        return url
    except Exception as e:
        return ""


async def _upload_to_supabase(file_path: str, doc_id: str, page_idx: int) -> str:
    """Upload rendered page image to Supabase Storage."""
    try:
        supabase = get_supabase_client()
        bucket_name = "documents"
        file_name = f"{doc_id}/pages/page_{page_idx}.png"
        
        with open(file_path, "rb") as f:
            file_data = f.read()
        
        # Upload file
        supabase.storage.from_(bucket_name).upload(
            file_name, file_data, {"content-type": "image/png", "upsert": "true"}
        )
        
        # Get public URL
        url = supabase.storage.from_(bucket_name).get_public_url(file_name)
        return url
    except Exception as e:
        return ""


async def _process_docx(file_path: str, flagged_findings: dict) -> dict:
    """Process DOCX file and generate heatmap data."""
    doc = Document(file_path)
    pages = []
    
    # Extract document_id from findings if available
    doc_id = flagged_findings.get("document_id", "temp")
    
    # DOCX doesn't have explicit pages, group paragraphs as "pages"
    paragraphs_per_page = 20
    page_num = 1
    current_page_paragraphs = []
    
    for para_idx, paragraph in enumerate(doc.paragraphs):
        current_page_paragraphs.append(paragraph)
        
        if len(current_page_paragraphs) >= paragraphs_per_page or para_idx == len(doc.paragraphs) - 1:
            boxes = []
            
            for para in current_page_paragraphs:
                text_content = para.text.strip()
                
                if text_content:
                    intensity, reason, source = _match_finding(
                        page_num, text_content, flagged_findings
                    )
                    
                    if intensity:
                        boxes.append({
                            "x": 0,
                            "y": 0,
                            "width": 100,
                            "height": 20,
                            "intensity": intensity,
                            "reason": reason,
                            "source_module": source
                        })
            
            image_url = await _render_docx_page(current_page_paragraphs, page_num - 1, doc_id)
            
            pages.append({
                "page_number": page_num,
                "image_url": image_url,
                "boxes": boxes
            })
            
            current_page_paragraphs = []
            page_num += 1
    
    return {"pages": pages}


async def _render_docx_page(paragraphs: list, page_idx: int, doc_id: str) -> str:
    """Render DOCX page as PNG and upload to Supabase."""
    try:
        img = Image.new("RGB", (800, 1000), "white")
        draw = ImageDraw.Draw(img)
        
        y_pos = 20
        for para in paragraphs:
            text = para.text.strip()[:80]
            if text:
                draw.text((20, y_pos), text, fill="black")
                y_pos += 40
        
        with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as tmp:
            img.save(tmp.name, "PNG")
            tmp_path = tmp.name
        
        url = await _upload_to_supabase(tmp_path, doc_id, page_idx)
        os.unlink(tmp_path)
        
        return url
    except Exception as e:
        return ""

from collections import Counter, defaultdict
from typing import Dict, List, Any, Tuple
from pathlib import Path
import statistics

from pptx import Presentation
from pptx.util import Pt, Emu
import fitz  # PyMuPDF
from docx import Document


def analyze_font_fingerprint(file_path: str, file_type: str) -> dict:
    """
    Analyze font consistency and detect formatting anomalies.
    
    Args:
        file_path: Path to the file
        file_type: 'pptx', 'pdf', or 'docx'
    
    Returns:
        Dict with consistency analysis, dominant fonts, and anomalies
    """
    if file_type == "pptx":
        return _analyze_pptx_fonts(file_path)
    elif file_type == "pdf":
        return _analyze_pdf_fonts(file_path)
    elif file_type == "docx":
        return _analyze_docx_fonts(file_path)
    else:
        return {
            "overall_consistency": "Unknown",
            "confidence": 0.0,
            "dominant_fonts": {},
            "font_anomalies": [],
            "spacing_anomalies": []
        }


def _analyze_pptx_fonts(file_path: str) -> dict:
    """Analyze font fingerprint for PPTX files."""
    prs = Presentation(file_path)
    dominant_fonts = {}
    font_anomalies = []
    spacing_anomalies = []
    
    for slide_idx, slide in enumerate(prs.slides):
        slide_num = slide_idx + 1
        slide_fonts = []
        slide_spacings = []
        
        # Extract all font runs and spacing from this slide
        for shape in slide.shapes:
            if not hasattr(shape, "text_frame"):
                continue
            
            for para_idx, paragraph in enumerate(shape.text_frame.paragraphs):
                # Collect spacing info
                if paragraph.space_before is not None:
                    slide_spacings.append(("before", int(paragraph.space_before)))
                if paragraph.space_after is not None:
                    slide_spacings.append(("after", int(paragraph.space_after)))
                
                para_fonts = []
                for run in paragraph.runs:
                    font_info = _extract_font_info_pptx(run)
                    if font_info:
                        slide_fonts.append(font_info)
                        para_fonts.append(font_info)
                
                # Detect font mixing within a single paragraph
                if len(para_fonts) > 1:
                    para_font_families = [f["family"] for f in para_fonts]
                    if len(set(para_font_families)) > 1:
                        font_anomalies.append({
                            "location": f"Slide {slide_num}, paragraph {para_idx + 1}",
                            "expected_font": para_font_families[0],
                            "found_font": ", ".join(set(para_font_families[1:])),
                            "severity": "medium"
                        })
        
        # Determine dominant font for this slide
        if slide_fonts:
            font_counter = Counter([f["signature"] for f in slide_fonts])
            dominant_sig = font_counter.most_common(1)[0][0]
            dominant_fonts[f"slide_{slide_num}"] = dominant_sig
            
            # Flag runs that deviate from dominant
            for font_info in slide_fonts:
                if font_info["signature"] != dominant_sig and font_counter[font_info["signature"]] < 3:
                    # Only flag if it's rare (appears less than 3 times)
                    if not any(a.get("found_font") == font_info["family"] for a in font_anomalies):
                        font_anomalies.append({
                            "location": f"Slide {slide_num}",
                            "expected_font": dominant_sig,
                            "found_font": font_info["signature"],
                            "severity": "low"
                        })
        
        # Detect spacing anomalies
        if len(slide_spacings) > 2:
            spacing_values = [v for _, v in slide_spacings]
            median_spacing = statistics.median(spacing_values)
            
            for spacing_type, value in slide_spacings:
                if median_spacing > 0 and value > median_spacing * 2:
                    spacing_anomalies.append({
                        "location": f"Slide {slide_num}",
                        "expected_spacing": f"{int(median_spacing / 12700)}pt",
                        "found_spacing": f"{int(value / 12700)}pt",
                        "severity": "low"
                    })
                    break  # Only report once per slide
    
    # Calculate overall consistency
    overall_consistency, confidence = _calculate_consistency(
        dominant_fonts, font_anomalies, spacing_anomalies
    )
    
    return {
        "overall_consistency": overall_consistency,
        "confidence": confidence,
        "dominant_fonts": dominant_fonts,
        "font_anomalies": font_anomalies[:10],  # Limit to top 10
        "spacing_anomalies": spacing_anomalies[:10]
    }


def _analyze_pdf_fonts(file_path: str) -> dict:
    """Analyze font fingerprint for PDF files with structural role awareness."""
    doc = fitz.open(file_path)
    dominant_fonts = {}
    font_anomalies = []
    spacing_anomalies = []
    
    for page_idx in range(len(doc)):
        page = doc[page_idx]
        page_num = page_idx + 1
        page_fonts = []
        page_spans = []  # Store spans with position for structural analysis
        
        blocks = page.get_text("dict")["blocks"]
        
        # First pass: collect all font info with positions
        for block in blocks:
            if block.get("type") != 0:  # Skip non-text blocks
                continue
            
            bbox = block.get("bbox", [0, 0, 0, 0])
            for line in block.get("lines", []):
                for span in line.get("spans", []):
                    font_info = _extract_font_info_pdf(span)
                    if font_info:
                        font_info["y_position"] = bbox[1]  # Top coordinate
                        font_info["text_length"] = len(span.get("text", ""))
                        page_fonts.append(font_info)
                        page_spans.append(font_info)
        
        if not page_fonts:
            continue
        
        # Build frequency distribution and classify structural roles
        font_counter = Counter([f["signature"] for f in page_fonts])
        size_counter = Counter([f["size"] for f in page_fonts])
        total_spans = len(page_fonts)
        
        # Calculate median size for role classification
        sizes = [f["size"] for f in page_fonts]
        median_size = statistics.median(sizes) if sizes else 12
        
        # Classify each span into structural roles
        structural_roles = _classify_structural_roles(page_spans, median_size)
        
        # Determine dominant font (most frequent)
        dominant_sig = font_counter.most_common(1)[0][0] if font_counter else "Unknown"
        dominant_fonts[f"page_{page_num}"] = dominant_sig
        
        # Only flag truly anomalous fonts (< 3% frequency threshold)
        frequency_threshold = max(1, int(total_spans * 0.03))
        
        anomalies_on_page = []
        for font_sig, count in font_counter.items():
            if font_sig != dominant_sig and count < frequency_threshold:
                # Check if this font is appropriate for its structural role
                spans_with_sig = [s for s in page_spans if s["signature"] == font_sig]
                
                # Don't flag if it's used consistently for a structural role
                roles = [structural_roles.get(id(s), "body") for s in spans_with_sig]
                if len(set(roles)) == 1 and roles[0] in ["header", "footer", "table_label"]:
                    continue  # Expected variation for structural role
                
                anomalies_on_page.append({
                    "signature": font_sig,
                    "count": count,
                    "frequency": count / total_spans
                })
        
        # Cap at top 5 most deviant per page
        anomalies_on_page.sort(key=lambda x: x["frequency"])
        for anomaly in anomalies_on_page[:5]:
            font_anomalies.append({
                "location": f"Page {page_num}",
                "expected_font": dominant_sig,
                "found_font": anomaly["signature"],
                "frequency": f"{anomaly['frequency']*100:.1f}%",
                "severity": "low"
            })
    
    doc.close()
    
    # Calculate overall consistency
    overall_consistency, confidence = _calculate_consistency(
        dominant_fonts, font_anomalies, spacing_anomalies
    )
    
    return {
        "overall_consistency": overall_consistency,
        "confidence": confidence,
        "dominant_fonts": dominant_fonts,
        "font_anomalies": font_anomalies[:10],
        "spacing_anomalies": spacing_anomalies[:10]
    }


def _classify_structural_roles(spans: list, median_size: float) -> dict:
    """Classify text spans into structural roles based on font size and position.
    
    Args:
        spans: List of font info dicts with y_position and size
        median_size: Median font size for the page
    
    Returns:
        Dict mapping span id() to role: "header", "footer", "table_label", "body"
    """
    roles = {}
    
    if not spans:
        return roles
    
    # Sort by y_position to find top/bottom regions
    sorted_spans = sorted(spans, key=lambda s: s.get("y_position", 0))
    
    # Top 10% and bottom 10% of page are likely headers/footers
    top_threshold = sorted_spans[max(0, len(sorted_spans) // 10)].get("y_position", 0)
    bottom_threshold = sorted_spans[max(0, -len(sorted_spans) // 10)].get("y_position", 999)
    
    for span in spans:
        y_pos = span.get("y_position", 0)
        size = span.get("size", median_size)
        is_bold = span.get("bold", False)
        
        # Header: larger than median OR bold OR in top region
        if size > median_size * 1.3 or (is_bold and size >= median_size * 0.9) or y_pos < top_threshold:
            roles[id(span)] = "header"
        # Footer: in bottom region or very small
        elif y_pos > bottom_threshold or size < median_size * 0.8:
            roles[id(span)] = "footer"
        # Table label: bold + short text + close to median size
        elif is_bold and span.get("text_length", 0) < 20:
            roles[id(span)] = "table_label"
        else:
            roles[id(span)] = "body"
    
    return roles


def _analyze_docx_fonts(file_path: str) -> dict:
    """Analyze font fingerprint for DOCX files."""
    doc = Document(file_path)
    dominant_fonts = {}
    font_anomalies = []
    spacing_anomalies = []
    
    # Group paragraphs by section (simplified: just sequential blocks)
    section_fonts = []
    section_spacings = []
    section_num = 1
    
    for para_idx, paragraph in enumerate(doc.paragraphs):
        # Collect spacing info
        if paragraph.paragraph_format.space_before is not None:
            section_spacings.append(("before", int(paragraph.paragraph_format.space_before)))
        if paragraph.paragraph_format.space_after is not None:
            section_spacings.append(("after", int(paragraph.paragraph_format.space_after)))
        
        para_fonts = []
        for run in paragraph.runs:
            font_info = _extract_font_info_docx(run)
            if font_info:
                section_fonts.append(font_info)
                para_fonts.append(font_info)
        
        # Detect font mixing within a paragraph
        if len(para_fonts) > 1:
            para_font_families = [f["family"] for f in para_fonts]
            if len(set(para_font_families)) > 1:
                font_anomalies.append({
                    "location": f"Paragraph {para_idx + 1}",
                    "expected_font": para_font_families[0],
                    "found_font": ", ".join(set(para_font_families[1:])),
                    "severity": "medium"
                })
        
        # Every 10 paragraphs, analyze as a section
        if (para_idx + 1) % 10 == 0 or para_idx == len(doc.paragraphs) - 1:
            if section_fonts:
                font_counter = Counter([f["signature"] for f in section_fonts])
                dominant_sig = font_counter.most_common(1)[0][0]
                dominant_fonts[f"section_{section_num}"] = dominant_sig
                section_num += 1
                section_fonts = []
            
            # Check spacing anomalies
            if len(section_spacings) > 2:
                spacing_values = [v for _, v in section_spacings]
                median_spacing = statistics.median(spacing_values)
                
                for spacing_type, value in section_spacings:
                    if median_spacing > 0 and value > median_spacing * 2:
                        spacing_anomalies.append({
                            "location": f"Section {section_num - 1}",
                            "expected_spacing": f"{int(median_spacing / 12700)}pt",
                            "found_spacing": f"{int(value / 12700)}pt",
                            "severity": "low"
                        })
                        break
                
                section_spacings = []
    
    # Calculate overall consistency
    overall_consistency, confidence = _calculate_consistency(
        dominant_fonts, font_anomalies, spacing_anomalies
    )
    
    return {
        "overall_consistency": overall_consistency,
        "confidence": confidence,
        "dominant_fonts": dominant_fonts,
        "font_anomalies": font_anomalies[:10],
        "spacing_anomalies": spacing_anomalies[:10]
    }


def _extract_font_info_pptx(run) -> dict:
    """Extract font information from a PPTX run."""
    try:
        font = run.font
        family = font.name if font.name else "Default"
        size = int(font.size.pt) if font.size else 12
        bold = font.bold if font.bold is not None else False
        italic = font.italic if font.italic is not None else False
        
        signature = f"{family} {size}pt"
        if bold:
            signature += " Bold"
        if italic:
            signature += " Italic"
        
        return {
            "family": family,
            "size": size,
            "bold": bold,
            "italic": italic,
            "signature": signature
        }
    except:
        return None


def _extract_font_info_pdf(span: dict) -> dict:
    """Extract font information from a PDF span."""
    try:
        font_name = span.get("font", "Unknown")
        size = int(span.get("size", 12))
        flags = span.get("flags", 0)
        
        # Parse flags: bit 0 = superscript, bit 1 = italic, bit 2 = serifed, bit 4 = bold
        bold = bool(flags & 16)
        italic = bool(flags & 2)
        
        # Clean font name (remove subset prefix like "ABCDEF+")
        if "+" in font_name:
            font_name = font_name.split("+", 1)[1]
        
        signature = f"{font_name} {size}pt"
        if bold:
            signature += " Bold"
        if italic:
            signature += " Italic"
        
        return {
            "family": font_name,
            "size": size,
            "bold": bold,
            "italic": italic,
            "signature": signature
        }
    except:
        return None


def _extract_font_info_docx(run) -> dict:
    """Extract font information from a DOCX run."""
    try:
        font = run.font
        family = font.name if font.name else "Default"
        size = int(font.size.pt) if font.size else 12
        bold = font.bold if font.bold is not None else False
        italic = font.italic if font.italic is not None else False
        
        signature = f"{family} {size}pt"
        if bold:
            signature += " Bold"
        if italic:
            signature += " Italic"
        
        return {
            "family": family,
            "size": size,
            "bold": bold,
            "italic": italic,
            "signature": signature
        }
    except:
        return None


def _calculate_consistency(
    dominant_fonts: dict, font_anomalies: list, spacing_anomalies: list
) -> Tuple[str, float]:
    """
    Calculate overall consistency rating and confidence.
    
    Returns:
        Tuple of (consistency_level, confidence_score)
    """
    total_anomalies = len(font_anomalies) + len(spacing_anomalies)
    pages_count = len(dominant_fonts)
    
    if pages_count == 0:
        return "Unknown", 0.0
    
    # Check if dominant fonts are consistent across pages/slides
    dominant_values = list(dominant_fonts.values())
    unique_dominants = len(set(dominant_values))
    
    # Calculate anomaly ratio
    anomaly_ratio = total_anomalies / max(pages_count, 1)
    
    # Determine consistency level
    if anomaly_ratio > 2.0 or unique_dominants > pages_count * 0.7:
        consistency = "Inconsistent"
        confidence = 0.85
    elif anomaly_ratio > 0.5 or unique_dominants > pages_count * 0.3:
        consistency = "Mixed"
        confidence = 0.75
    else:
        consistency = "Consistent"
        confidence = 0.90
    
    return consistency, confidence

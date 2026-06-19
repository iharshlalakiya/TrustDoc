import re
from pathlib import Path

import pdfplumber
from docx import Document as DocxDocument
from pptx import Presentation

BUSINESS_DOMAIN_KEYWORDS: dict[str, list[str]] = {
    "AI/ML": [
        "machine learning",
        "deep learning",
        "neural network",
        "artificial intelligence",
        "llm",
        "gpt",
        "transformer",
        "computer vision",
        "nlp",
        "natural language",
        "model training",
        "inference",
        "dataset",
        "classification",
        "regression",
        "tensorflow",
        "pytorch",
    ],
    "Cybersecurity": [
        "cybersecurity",
        "encryption",
        "firewall",
        "malware",
        "ransomware",
        "phishing",
        "vulnerability",
        "penetration test",
        "zero trust",
        "authentication",
        "authorization",
        "incident response",
        "threat",
        "soc",
        "siem",
        "intrusion",
    ],
    "Finance": [
        "revenue",
        "profit",
        "investment",
        "portfolio",
        "financial",
        "accounting",
        "balance sheet",
        "cash flow",
        "roi",
        "valuation",
        "equity",
        "dividend",
        "budget",
        "forecast",
        "fiscal",
        "audit",
    ],
    "Legal": [
        "contract",
        "agreement",
        "clause",
        "liability",
        "jurisdiction",
        "compliance",
        "regulation",
        "litigation",
        "arbitration",
        "intellectual property",
        "confidential",
        "terms and conditions",
        "warranty",
        "indemnity",
        "legal counsel",
    ],
    "Healthcare": [
        "patient",
        "clinical",
        "diagnosis",
        "treatment",
        "hospital",
        "medical",
        "healthcare",
        "pharmaceutical",
        "fda",
        "hipaa",
        "symptoms",
        "therapy",
        "physician",
        "nurse",
        "ehr",
    ],
}

DOCUMENT_TYPE_KEYWORDS: dict[str, list[str]] = {
    "pitch deck": [
        "pitch",
        "investor",
        "fundraising",
        "series a",
        "seed round",
        "go-to-market",
        "traction",
        "problem",
        "solution",
        "market size",
    ],
    "report": [
        "executive summary",
        "findings",
        "methodology",
        "analysis",
        "conclusion",
        "recommendation",
        "overview",
        "assessment",
    ],
    "resume": [
        "experience",
        "education",
        "skills",
        "curriculum vitae",
        "references",
        "employment history",
        "professional summary",
    ],
    "invoice": [
        "invoice",
        "bill to",
        "due date",
        "payment terms",
        "subtotal",
        "tax",
        "amount due",
        "invoice number",
    ],
    "contract": [
        "party",
        "whereas",
        "hereby",
        "agreement",
        "obligations",
        "termination",
        "governing law",
        "signature",
    ],
    "presentation": [
        "agenda",
        "slide",
        "overview",
        "thank you",
        "questions",
        "objectives",
        "roadmap",
    ],
}

TECH_TERMS: list[str] = [
    "tensorflow",
    "pytorch",
    "blockchain",
    "ipfs",
    "kubernetes",
    "docker",
    "aws",
    "azure",
    "gcp",
    "redis",
    "postgresql",
    "mongodb",
    "react",
    "fastapi",
    "gemini",
    "openai",
    "huggingface",
    "transformers",
    "opencv",
    "ethereum",
    "smart contract",
    "microservices",
    "api",
    "serverless",
    "nginx",
    "kafka",
    "spark",
    "hadoop",
    "snowflake",
    "databricks",
]

INFRA_TERMS: list[str] = [
    "aws",
    "azure",
    "gcp",
    "cloud",
    "kubernetes",
    "docker",
    "terraform",
    "ansible",
    "ci/cd",
    "pipeline",
    "load balancer",
    "vpc",
    "subnet",
    "firewall",
    "cdn",
    "s3",
    "ec2",
    "lambda",
    "server",
    "datacenter",
    "on-premise",
    "hybrid cloud",
]

NAME_PATTERN = re.compile(
    r"\b(?:Mr\.|Ms\.|Mrs\.|Dr\.)?\s*([A-Z][a-z]+(?:\s+[A-Z][a-z]+){1,2})\b"
)
STAT_PATTERN = re.compile(
    r"\b(?:\$|€|£)?\d[\d,]*(?:\.\d+)?(?:%|\s*(?:million|billion|k|m|b))?\b",
    re.IGNORECASE,
)


def _normalize_text(text: str) -> str:
    return re.sub(r"\s+", " ", text).strip().lower()


def _extract_pptx(file_path: Path) -> tuple[str, int, int]:
    presentation = Presentation(str(file_path))
    slide_count = len(presentation.slides)
    hidden_slides = 0

    text_parts: list[str] = []
    for slide in presentation.slides:
        if slide._element.get("show") == "0":
            hidden_slides += 1

        for shape in slide.shapes:
            if shape.has_text_frame:
                text_parts.append(shape.text_frame.text)

    return "\n".join(text_parts), slide_count, hidden_slides


def _extract_pdf(file_path: Path) -> tuple[str, int]:
    text_parts: list[str] = []
    page_count = 0

    with pdfplumber.open(str(file_path)) as pdf:
        page_count = len(pdf.pages)
        for page in pdf.pages:
            page_text = page.extract_text()
            if page_text:
                text_parts.append(page_text)

    return "\n".join(text_parts), page_count


def _extract_docx(file_path: Path) -> tuple[str, int]:
    document = DocxDocument(str(file_path))
    paragraphs = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
    page_count = max(1, len(paragraphs) // 30)
    return "\n".join(paragraphs), page_count


def _extract_content(file_path: Path, file_type: str) -> tuple[str, int, int]:
    normalized_type = file_type.lower().lstrip(".")

    if normalized_type == "pptx":
        text, count, hidden = _extract_pptx(file_path)
        return text, count, hidden
    if normalized_type == "pdf":
        text, count = _extract_pdf(file_path)
        return text, count, 0
    if normalized_type == "docx":
        text, count = _extract_docx(file_path)
        return text, count, 0

    raise ValueError(f"Unsupported file type: {file_type}")


def _score_keywords(text: str, keywords: dict[str, list[str]]) -> tuple[str, int]:
    best_label = "General"
    best_score = 0

    for label, terms in keywords.items():
        score = sum(1 for term in terms if term in text)
        if score > best_score:
            best_score = score
            best_label = label

    return best_label, best_score


def _detect_business_domain(text: str) -> str:
    domain, score = _score_keywords(text, BUSINESS_DOMAIN_KEYWORDS)
    return domain if score > 0 else "General"


def _detect_document_type(text: str, file_type: str) -> str:
    doc_type, score = _score_keywords(text, DOCUMENT_TYPE_KEYWORDS)

    if score == 0:
        if file_type.lower() == "pptx":
            return "presentation"
        return "report"

    return doc_type


NAME_BLOCKLIST = {
    "pitch",
    "deck",
    "report",
    "summary",
    "overview",
    "agenda",
    "introduction",
    "conclusion",
    "machine",
    "learning",
}


def _extract_names(text: str) -> list[str]:
    names: list[str] = []
    for match in NAME_PATTERN.finditer(text):
        name = match.group(1).strip()
        words = name.split()
        if "\n" in name:
            continue
        if any(word.lower() in NAME_BLOCKLIST for word in words):
            continue
        if name.lower() not in {"the", "and", "for", "with"}:
            names.append(name)
    return list(dict.fromkeys(names))[:10]


def _extract_tech_terms(text: str) -> list[str]:
    normalized = _normalize_text(text)
    found = [term for term in TECH_TERMS if term in normalized]
    return list(dict.fromkeys(found))[:15]


def _extract_statistics(text: str) -> list[str]:
    matches = STAT_PATTERN.findall(text)
    return list(dict.fromkeys(matches))[:15]


def _extract_infra_terms(text: str) -> list[str]:
    normalized = _normalize_text(text)
    found = [term for term in INFRA_TERMS if term in normalized]
    return list(dict.fromkeys(found))[:15]


def classify_document(file_path: str, file_type: str) -> dict:
    """
    Classify a document by format, domain, type, and extracted entities.

    Args:
        file_path: Absolute or relative path to the document on disk.
        file_type: File extension without dot (pptx, pdf, docx).

    Returns:
        Classification result with metadata and key entities.
    """
    path = Path(file_path)
    if not path.is_file():
        raise FileNotFoundError(f"Document not found: {file_path}")

    normalized_type = file_type.lower().lstrip(".")
    raw_text, slide_or_page_count, hidden_slides = _extract_content(path, normalized_type)
    normalized_text = _normalize_text(raw_text)

    business_domain = _detect_business_domain(normalized_text)
    document_type = _detect_document_type(normalized_text, normalized_type)

    result = {
        "document_type": document_type,
        "file_format": normalized_type,
        "slide_or_page_count": slide_or_page_count,
        "hidden_slides": hidden_slides if normalized_type == "pptx" else None,
        "business_domain": business_domain,
        "key_entities": {
            "team": _extract_names(raw_text),
            "tech": _extract_tech_terms(normalized_text),
            "data": _extract_statistics(raw_text),
            "infra": _extract_infra_terms(normalized_text),
        },
        "confidence": 0.90,
    }

    return result

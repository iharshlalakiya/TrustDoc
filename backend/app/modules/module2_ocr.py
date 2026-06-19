import re
import unicodedata
from pathlib import Path

import fitz
from docx import Document as DocxDocument
from pptx import Presentation

ZERO_WIDTH_CHARS = {"\u200b", "\u200c", "\u200d", "\ufeff"}
RTL_OVERRIDE_CHARS = {"\u202e", "\u202b"}
LOOKALIKE_PAIRS = (
    ("0", "O"),
    ("O", "0"),
    ("|", "I"),
    ("I", "|"),
    ("1", "l"),
    ("l", "1"),
)
SUSPICIOUS_UNICODE_CATEGORIES = {"Cf", "Co", "Cn", "Cs"}
DOCX_PARAGRAPHS_PER_PAGE = 30

LOOKALIKE_SUBSTITUTIONS = str.maketrans(
    {
        "0": "O",
        "O": "0",
        "|": "I",
        "I": "|",
        "1": "l",
        "l": "1",
    }
)


def _extract_segments(file_path: Path, file_type: str) -> list[tuple[int, str]]:
    normalized_type = file_type.lower().lstrip(".")

    if normalized_type == "pdf":
        return _extract_pdf_segments(file_path)
    if normalized_type == "pptx":
        return _extract_pptx_segments(file_path)
    if normalized_type == "docx":
        return _extract_docx_segments(file_path)

    raise ValueError(f"Unsupported file type: {file_type}")


def _extract_pdf_segments(file_path: Path) -> list[tuple[int, str]]:
    segments: list[tuple[int, str]] = []
    with fitz.open(str(file_path)) as document:
        for page_number, page in enumerate(document, start=1):
            segments.append((page_number, page.get_text()))
    return segments


def _extract_pptx_segments(file_path: Path) -> list[tuple[int, str]]:
    presentation = Presentation(str(file_path))
    segments: list[tuple[int, str]] = []

    for slide_number, slide in enumerate(presentation.slides, start=1):
        text_parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_parts.append(shape.text_frame.text)
        segments.append((slide_number, "\n".join(text_parts)))

    return segments


def _extract_docx_segments(file_path: Path) -> list[tuple[int, str]]:
    document = DocxDocument(str(file_path))
    paragraphs = [paragraph.text for paragraph in document.paragraphs]

    if not paragraphs:
        return [(1, "")]

    segments: list[tuple[int, str]] = []
    for index in range(0, len(paragraphs), DOCX_PARAGRAPHS_PER_PAGE):
        page_number = index // DOCX_PARAGRAPHS_PER_PAGE + 1
        chunk = "\n".join(paragraphs[index : index + DOCX_PARAGRAPHS_PER_PAGE])
        segments.append((page_number, chunk))

    return segments


def _detect_zero_width(text: str) -> list[str]:
    found = sorted({char for char in text if char in ZERO_WIDTH_CHARS})
    return [f"Zero-width character U+{ord(char):04X}" for char in found]


def _detect_rtl_override(text: str) -> list[str]:
    found = sorted({char for char in text if char in RTL_OVERRIDE_CHARS})
    return [f"RTL override character U+{ord(char):04X}" for char in found]


def _detect_lookalike_substitutions(text: str) -> list[str]:
    anomalies: list[str] = []

    for token in re.findall(r"\b[\w|]+\b", text):
        if "|" in token:
            anomalies.append(f"Pipe character used in token: '{token}'")

        if "0" in token and "O" in token:
            anomalies.append(f"Mixed '0' and 'O' in token: '{token}'")

        if re.search(r"(?<=[a-zA-Z])1(?=[a-zA-Z])", token) or re.search(
            r"(?<=[a-zA-Z])l(?=\d)|(?<=\d)l(?=[a-zA-Z])", token
        ):
            anomalies.append(f"Possible '1'/'l' confusion in token: '{token}'")

    return anomalies


def _detect_unicode_anomalies(text: str) -> list[str]:
    anomalies: list[str] = []

    for index, char in enumerate(text):
        if char in ZERO_WIDTH_CHARS or char in RTL_OVERRIDE_CHARS:
            continue

        category = unicodedata.category(char)
        if category in SUSPICIOUS_UNICODE_CATEGORIES:
            anomalies.append(
                f"Suspicious Unicode {category} at position {index}: U+{ord(char):04X}"
            )
            continue

        if ord(char) > 127:
            name = unicodedata.name(char, "")
            if name.startswith(("CYRILLIC", "GREEK", "ARMENIAN")):
                context_start = max(0, index - 10)
                context_end = min(len(text), index + 10)
                snippet = text[context_start:context_end].replace("\n", " ")
                anomalies.append(
                    f"Non-ASCII homoglyph candidate U+{ord(char):04X} ({name}) near '{snippet}'"
                )

    return anomalies


def _build_homoglyph_pattern(term: str) -> re.Pattern[str]:
    char_map = {
        "o": "[oO0]",
        "i": "[iI|]",
        "l": "[l1]",
    }
    pattern_parts = []
    for char in term:
        lower = char.lower()
        if lower in char_map:
            pattern_parts.append(char_map[lower])
        elif char in {"0", "1"}:
            pattern_parts.append("[oO0]" if char == "0" else "[l1]")
        else:
            pattern_parts.append(re.escape(char))
    return re.compile(r"\b" + "".join(pattern_parts) + r"\b", re.IGNORECASE)


def _check_term_consistency(text: str, expected_terms: list[str]) -> list[str]:
    anomalies: list[str] = []
    lowered_text = text.lower()

    for term in expected_terms:
        normalized_term = term.strip().lower()
        if not normalized_term:
            continue

        if normalized_term in lowered_text:
            continue

        pattern = _build_homoglyph_pattern(normalized_term)
        for match in pattern.finditer(text):
            matched = match.group()
            if matched.lower() != normalized_term:
                anomalies.append(
                    f"Inconsistent spelling for '{term}': found '{matched}'"
                )

        if normalized_term not in lowered_text and not pattern.search(text):
            normalized_variant = normalized_term.translate(LOOKALIKE_SUBSTITUTIONS)
            if normalized_variant != normalized_term and normalized_variant in lowered_text:
                anomalies.append(
                    f"Expected term '{term}' appears with look-alike substitutions"
                )

    return anomalies


def _calculate_confidence(anomaly_count: int) -> float:
    if anomaly_count == 0:
        return 1.0
    confidence = max(0.0, 1.0 - (anomaly_count * 0.08))
    return round(confidence, 2)


def validate_ocr(
    file_path: str,
    file_type: str,
    expected_terms: list[str] | None = None,
) -> dict:
    """
    Validate extracted document text for OCR and homoglyph anomalies.

    Args:
        file_path: Path to the document on disk.
        file_type: File extension without dot (pptx, pdf, docx).
        expected_terms: Technical terms from module 1 to verify consistent spelling.

    Returns:
        OCR validation report with anomaly details and confidence score.
    """
    path = Path(file_path)
    if not path.is_file():
        raise FileNotFoundError(f"Document not found: {file_path}")

    segments = _extract_segments(path, file_type)
    expected_terms = expected_terms or []

    anomalies_found: list[str] = []
    suspicious_slides_or_pages: list[int] = []
    homoglyph_detected = False
    rtl_override_detected = False
    zero_width_chars_detected = False

    full_text_parts: list[str] = []

    for segment_index, segment_text in segments:
        full_text_parts.append(segment_text)
        segment_anomalies: list[str] = []

        zero_width = _detect_zero_width(segment_text)
        if zero_width:
            zero_width_chars_detected = True
            segment_anomalies.extend(zero_width)

        rtl = _detect_rtl_override(segment_text)
        if rtl:
            rtl_override_detected = True
            segment_anomalies.extend(rtl)

        lookalikes = _detect_lookalike_substitutions(segment_text)
        if lookalikes:
            homoglyph_detected = True
            segment_anomalies.extend(lookalikes)

        unicode_anomalies = _detect_unicode_anomalies(segment_text)
        if unicode_anomalies:
            homoglyph_detected = True
            segment_anomalies.extend(unicode_anomalies)

        if segment_anomalies:
            suspicious_slides_or_pages.append(segment_index)
            anomalies_found.extend(
                [f"[page/slide {segment_index}] {item}" for item in segment_anomalies]
            )

    full_text = "\n".join(full_text_parts)
    term_anomalies = _check_term_consistency(full_text, expected_terms)
    if term_anomalies:
        homoglyph_detected = True
        anomalies_found.extend(term_anomalies)

    anomalies_found = list(dict.fromkeys(anomalies_found))
    suspicious_slides_or_pages = sorted(set(suspicious_slides_or_pages))
    confidence = _calculate_confidence(len(anomalies_found))

    return {
        "text_extraction_clean": len(anomalies_found) == 0,
        "confidence": confidence,
        "anomalies_found": anomalies_found,
        "suspicious_slides_or_pages": suspicious_slides_or_pages,
        "homoglyph_detected": homoglyph_detected,
        "rtl_override_detected": rtl_override_detected,
        "zero_width_chars_detected": zero_width_chars_detected,
    }
